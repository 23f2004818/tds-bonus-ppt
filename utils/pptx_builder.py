from pptx import Presentation
from pptx.util import Inches, Pt

def build_pptx_from_template(template_file, slides, output_file="output.pptx"):
    prs = Presentation(template_file)
    blank_layout = prs.slide_layouts[1]  # Title + Content

    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title_placeholder.text = slide_data["title"]
        content_placeholder.text = slide_data["content"]

    prs.save(output_file)
    return output_file
